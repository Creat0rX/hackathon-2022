
### Settings ###
"""Get coronavirus case statistics from the world and USA states, store them in a database, and create a powerpoint with findings"""
import os
from glob import glob
from RPA.Browser.Selenium import Selenium
from RPA.Email.ImapSmtp import ImapSmtp
from RPA.Email.Exchange import Exchange
from RPA.Robocorp.Vault import Vault
from robot.api import logger
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
import sqlite3 as sqlite

### Variables ###
XPATH_MAIN = "//table[@id='main_table_countries_today']/tbody/tr[@role='row'][@class='odd' or @class='even']"
XPATH_USA = "//table[@id='usa_table_countries_today']/tbody[1]/tr[not(@class='total_row_usa odd' or @class='total_row')]"
XPATH_GRAPH = "//div[@class='col-md-12']/div"
COUNTRY_LIST = ["Country_link", "Rank", "Country", "Total_cases", "New_cases", "Total_deaths", "New_deaths",
                "Total_recovered", "New_recovered", "Active_cases", "Serious_cases", "Cases_per_mil",
                "Deaths_per_mil", "Total_tests", "Tests_per_mil", "Population"]
STATE_LIST = ["State_link", "Rank", "State", "Total_cases", "New_cases", "Total_deaths", "New_deaths",
                "Total_recovered", "Active_cases", "Cases_per_mil", "Deaths_per_mil", "Total_tests",
                "Tests_per_mil", "Population"]
SLASH = os.sep
driver = Selenium()
conn = None
main_table_countries = []
main_table_states = []

### Keywords ###
def delete_powerpoint_if_exists():
    log_info_message("delete_powerpoint_if_exists")
    for f in glob(f"output{SLASH}*.png"):
        os.remove(f)

def open_browser_to_covid_website():
    log_info_message("open_browser_to_covid_website")
    driver.open_available_browser("https://www.worldometers.info/coronavirus/")

def scrape_table_from_website():
    log_info_message("scrape_table_from_website")
    counter = driver.get_element_count(XPATH_MAIN)
    for n in range(counter):
        main_table_countries.append(insert_values_safely(n + 1, COUNTRY_LIST, XPATH_MAIN))

def scrape_us_table_from_website():
    log_info_message("scrape_us_table_from_website")
    driver.open_available_browser("https://www.worldometers.info/coronavirus/country/us/")
    counter = driver.get_element_count(XPATH_USA)
    for n in range(counter):
        main_table_states.append(insert_values_safely(n + 1, STATE_LIST, XPATH_USA))

def screenshot_us_graphs_from_website():
    log_info_message("screenshot_us_graphs_from_website")
    for state in main_table_states[:10]:
        driver.go_to(state["State_link"])
        driver.screenshot(f"{XPATH_GRAPH}[@class='tabbable-panel-cases']", f"output{SLASH}total_cases.png")
        driver.screenshot(f"{XPATH_GRAPH}[@id='graph-cases-daily']", f"output{SLASH}daily_new_cases.png")
        driver.screenshot(f"{XPATH_GRAPH}[@id='graph-active-cases-total']", f"output{SLASH}active_cases.png")
        driver.screenshot(f"{XPATH_GRAPH}[@class='tabbable-panel-deaths']", f"output{SLASH}total_deaths.png")
        add_to_powerpoint(state["State"])

def insert_values_safely(n, list_choice, link):
    table_values = {}

    for i, header in enumerate(list_choice):
        try:
            if i == 0:
                table_values[header] = driver.get_element_attribute(f"{link}[{n}]/td[2]/a", "href")
            else:
                table_values[header] = driver.get_element_attribute(f"{link}[{n}]/td[{i}]", "innerText")
        except:
            table_values[header] = ""
    return table_values

def connect_to_sql_database():
    log_info_message("connect_to_sql_database")
    conn = sqlite.connect(f"output{SLASH}corona.db")
    return conn

def add_sql_tables(conn):
    log_info_message("add_sql_tables")
    if conn is not None:
        cursor = conn.cursor()

        tables_exist = len(list(cursor.execute("SELECT name FROM sqlite_master WHERE type='table';").fetchall()))
        if not tables_exist:
            sql_bool = False
        else:
            sql_bool = True

        sql_countries = f"""CREATE TABLE IF NOT EXISTS countries(
                        Rank integer PRIMARY KEY, {' text,'.join(COUNTRY_LIST[2:])} text);"""
        sql_states = f"""CREATE TABLE IF NOT EXISTS american_states(
                        Rank integer PRIMARY KEY, {' text,'.join(STATE_LIST[2:])} text);"""

        cursor.execute(sql_countries)
        cursor.execute(sql_states)
        conn.commit()
        add_values_to_sql_tables(conn, sql_bool)
    else:
        print("Failed to establish a connection with the database.")

def add_values_to_sql_tables(conn, sql_bool):
    log_info_message("add_values_to_sql_tables")
    cursor = conn.cursor()
    if sql_bool:
        # UPDATE SQL IF TABLES EXISTED
        query1 = f"""UPDATE countries SET {' = ? , '.join(COUNTRY_LIST[2:])} = ? WHERE Rank = ?;"""
        query2 = f"""UPDATE american_states SET {' = ? , '.join(STATE_LIST[2:])} = ? WHERE Rank = ?;"""
        for country in main_table_countries:
            country_details = (country["Country"], country["Total_cases"], country["New_cases"], 
                country["Total_deaths"], country["New_deaths"], country["Total_recovered"], 
                country["New_recovered"], country["Active_cases"], country["Serious_cases"], 
                country["Cases_per_mil"], country["Deaths_per_mil"], country["Total_tests"], 
                country["Tests_per_mil"], country["Population"], int(country["Rank"]))
            cursor.execute(query1, country_details)
            conn.commit()
        for state in main_table_states:
            print(state)
            state_details = (state["State"], state["Total_cases"], state["New_cases"], state["Total_deaths"],
                state["New_deaths"], state["Total_recovered"], state["Active_cases"], state["Cases_per_mil"],
                state["Deaths_per_mil"], state["Total_tests"], state["Tests_per_mil"], state["Population"],
                int(state["Rank"]))
            cursor.execute(query2, state_details)
            conn.commit()
    else:
        # INSERT SQL IF TABLES NEVER EXISTED
        query1 = f"""INSERT INTO countries ({', '.join(COUNTRY_LIST[1:])})
                    VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?);"""
        query2 = f"""INSERT INTO american_states ({', '.join(STATE_LIST[1:])})
                    VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?);"""
        for country in main_table_countries:
            country_details = (int(country["Rank"]), country["Country"], country["Total_cases"], 
                country["New_cases"], country["Total_deaths"], country["New_deaths"],
                country["Total_recovered"], country["New_recovered"], country["Active_cases"],
                country["Serious_cases"], country["Cases_per_mil"], country["Deaths_per_mil"],
                country["Total_tests"], country["Tests_per_mil"], country["Population"])
            cursor.execute(query1, country_details)
            conn.commit()
        for state in main_table_states:
            state_details = (int(state["Rank"]), state["State"], state["Total_cases"], state["New_cases"],
                state["Total_deaths"], state["New_deaths"], state["Total_recovered"],
                state["Active_cases"], state["Cases_per_mil"], state["Deaths_per_mil"],
                state["Total_tests"], state["Tests_per_mil"], state["Population"])
            cursor.execute(query2, state_details)
            conn.commit()

def disconnect_from_sql_database():
    log_info_message("disconnect_from_sql_database")
    if conn is not None:
        conn.close()

def add_to_powerpoint(title):
    log_info_message("add_to_powerpoint")
    try:
        ppt = Presentation(f"output{SLASH}presentation-{str(datetime.now().date())}.pptx")
    except:
        ppt = Presentation()
    title_slide = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(title_slide)
    slide.shapes.title.text = f"USA - {title}"
    slide.shapes.add_picture(f"output{SLASH}total_cases.png", left=Pt(50), top=Pt(100), height=Pt(200), width=Pt(300))
    slide.shapes.add_picture(f"output{SLASH}active_cases.png", left=Pt(50), top=Pt(310), height=Pt(200), width=Pt(300))
    slide.shapes.add_picture(f"output{SLASH}daily_new_cases.png", left=Pt(400), top=Pt(100), height=Pt(200), width=Pt(300))
    slide.shapes.add_picture(f"output{SLASH}total_deaths.png", left=Pt(400), top=Pt(310), height=Pt(200), width=Pt(300))
    ppt.save(f"output{SLASH}presentation-{str(datetime.now().date())}.pptx")

def delete_screenshots():
    log_info_message("delete_screenshots")
    for f in glob(f"output{SLASH}*.png"):
        os.remove(f)

def get_credentials():
    log_info_message("get_credentials")
    _credentials = Vault().get_secret("hackathon2022")
    return _credentials["credentials"]

def send_email():
    log_info_message("send_email")
    _credentials = get_credentials()
    email = ImapSmtp()
    email.authorize_smtp(_credentials["user"], _credentials["pword"], "smtp.gmail.com", 587)
    email.send_smtp_hello()
    sender = _credentials["user"]
    recipients = _credentials["receiver"]
    subject = "Hackathon 2022"
    body = "Please see attachment for powerpoint presentation. The attached database uses sqlite and \
        so does not use login credentials."
    attachments = [f"output{SLASH}presentation-{str(datetime.now().date())}.pptx", f"output{SLASH}corona.db"]
    email.send_message(sender, recipients, subject, body, attachments)

def close_browser():
    log_info_message("close_browser")
    if driver is not None:
        driver.close_browser()

def log_info_message(message):
    logger.info(message, True, True)

def main():
    global conn
    try:
        open_browser_to_covid_website()
        scrape_table_from_website()
        scrape_us_table_from_website()
        conn = connect_to_sql_database()
        add_sql_tables(conn)
        delete_powerpoint_if_exists()
        screenshot_us_graphs_from_website()
        delete_screenshots()
        send_email()
    except Exception as e:
        logger.error(f"An error occurred: {str(e)}")
    finally:
        disconnect_from_sql_database()
        close_browser()

### Tasks ###
if __name__ == "__main__":
    log_info_message(__doc__)
    main()
